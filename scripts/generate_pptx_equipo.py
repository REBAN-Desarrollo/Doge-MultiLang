#!/usr/bin/env python3
"""
generate_pptx_equipo.py
-----------------------
Generates a professional PPTX presentation for the QPH Multi-Idioma
non-technical team meeting.

Output: QPH_Problematicas_Equipo.pptx

Usage:
    python scripts/generate_pptx_equipo.py
"""

import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# CONSTANTS
# ---------------------------------------------------------------------------
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Colors
DARK_BLUE = RGBColor(0x1B, 0x2A, 0x4A)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MEDIUM_GRAY = RGBColor(0xDD, 0xDD, 0xDD)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xC0, 0x39, 0x2B)
LIGHT_BLUE = RGBColor(0x2C, 0x3E, 0x6D)
SOFT_ORANGE_BG = RGBColor(0xFD, 0xF2, 0xE9)

FONT_NAME = "Calibri"

# Output path
SCRIPT_DIR = Path(__file__).resolve().parent
PROJECT_ROOT = SCRIPT_DIR.parent
OUTPUT_PATH = PROJECT_ROOT / "QPH_Problematicas_Equipo.pptx"


# ---------------------------------------------------------------------------
# HELPER FUNCTIONS
# ---------------------------------------------------------------------------

def set_slide_bg(slide, color):
    """Set the background color of a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color):
    """Add a filled rectangle (no border) as a background block."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def set_text(tf, text, font_size=Pt(18), bold=False, color=BLACK,
             alignment=PP_ALIGN.LEFT, font_name=FONT_NAME, space_after=Pt(6)):
    """Set text in a text frame, clearing existing content."""
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = space_after


def add_paragraph(tf, text, font_size=Pt(18), bold=False, color=BLACK,
                  alignment=PP_ALIGN.LEFT, font_name=FONT_NAME,
                  space_before=Pt(0), space_after=Pt(6), level=0):
    """Add a new paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    p.level = level
    return p


def add_textbox(slide, left, top, width, height, text, font_size=Pt(18),
                bold=False, color=BLACK, alignment=PP_ALIGN.LEFT,
                font_name=FONT_NAME, word_wrap=True, anchor=MSO_ANCHOR.TOP):
    """Add a text box to the slide and return the shape."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    txBox.text_frame.auto_size = None
    # Set vertical anchor
    try:
        tf.paragraphs[0].alignment = alignment
    except Exception:
        pass

    set_text(tf, text, font_size=font_size, bold=bold, color=color,
             alignment=alignment, font_name=font_name)
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=Pt(18),
                    color=DARK_GRAY, bold_items=None, highlight_items=None,
                    font_name=FONT_NAME):
    """
    Add a bullet list to the slide.
    bold_items: set of indices to bold
    highlight_items: set of indices to color with ORANGE
    """
    if bold_items is None:
        bold_items = set()
    if highlight_items is None:
        highlight_items = set()

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.font.size = font_size
        p.font.name = font_name
        p.space_after = Pt(8)

        item_color = ORANGE if i in highlight_items else color
        is_bold = i in bold_items

        # Handle mixed formatting with runs
        if isinstance(item, tuple):
            # Tuple: (prefix, bold_part, suffix)
            prefix, bold_part, suffix = item
            if prefix:
                run = p.add_run()
                run.text = prefix
                run.font.size = font_size
                run.font.name = font_name
                run.font.color.rgb = item_color
                run.font.bold = False
            run_b = p.add_run()
            run_b.text = bold_part
            run_b.font.size = font_size
            run_b.font.name = font_name
            run_b.font.color.rgb = item_color
            run_b.font.bold = True
            if suffix:
                run_s = p.add_run()
                run_s.text = suffix
                run_s.font.size = font_size
                run_s.font.name = font_name
                run_s.font.color.rgb = item_color
                run_s.font.bold = False
        else:
            p.text = item
            p.font.color.rgb = item_color
            p.font.bold = is_bold

    return txBox


def build_table(slide, left, top, width, height, rows, col_widths=None,
                header_color=DARK_BLUE, header_font_color=WHITE,
                font_size=Pt(12)):
    """
    Build a table on the slide.
    rows: list of lists (first row = header).
    col_widths: list of Inches values (optional).
    Returns the table shape.
    """
    num_rows = len(rows)
    num_cols = len(rows[0]) if rows else 0

    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table

    # Set column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r_idx, row in enumerate(rows):
        for c_idx, cell_text in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = ""

            # Set cell text
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = str(cell_text)
            p.font.size = font_size
            p.font.name = FONT_NAME

            if r_idx == 0:
                # Header row
                p.font.bold = True
                p.font.color.rgb = header_font_color
                p.alignment = PP_ALIGN.CENTER
                cell_fill = cell.fill
                cell_fill.solid()
                cell_fill.fore_color.rgb = header_color
            else:
                p.font.color.rgb = DARK_GRAY
                p.alignment = PP_ALIGN.LEFT
                # Alternating row colors
                cell_fill = cell.fill
                cell_fill.solid()
                if r_idx % 2 == 0:
                    cell_fill.fore_color.rgb = RGBColor(0xEB, 0xEB, 0xEB)
                else:
                    cell_fill.fore_color.rgb = WHITE

            # Vertical centering
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return table_shape


# ---------------------------------------------------------------------------
# SLIDE HEADER HELPER
# ---------------------------------------------------------------------------
def add_slide_header(slide, title_text, subtitle_text=None):
    """Add a dark blue header bar at the top of a content slide."""
    # Header background bar
    add_shape_bg(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.3), DARK_BLUE)

    # Title text
    add_textbox(slide, Inches(0.7), Inches(0.2), Inches(11.5), Inches(0.7),
                title_text, font_size=Pt(30), bold=True, color=WHITE,
                alignment=PP_ALIGN.LEFT)

    if subtitle_text:
        add_textbox(slide, Inches(0.7), Inches(0.8), Inches(11.5), Inches(0.4),
                    subtitle_text, font_size=Pt(16), bold=False, color=ORANGE,
                    alignment=PP_ALIGN.LEFT)


def add_slide_number(slide, number, total=26):
    """Add a small slide number in the bottom right."""
    add_textbox(slide, Inches(12.0), Inches(7.0), Inches(1.0), Inches(0.4),
                f"{number}/{total}", font_size=Pt(10), bold=False,
                color=RGBColor(0x99, 0x99, 0x99), alignment=PP_ALIGN.RIGHT)


def add_highlight_box(slide, left, top, width, height, text, font_size=Pt(16)):
    """Add an orange-highlighted callout box."""
    shape = add_shape_bg(slide, left, top, width, height, ORANGE)
    shape.text_frame.word_wrap = True
    set_text(shape.text_frame, text, font_size=font_size, bold=True,
             color=WHITE, alignment=PP_ALIGN.CENTER)
    return shape


def add_quote_box(slide, left, top, width, height, text, font_size=Pt(16)):
    """Add a quote-style box with left orange border feel."""
    # Background
    bg = add_shape_bg(slide, left, top, width, height, SOFT_ORANGE_BG)
    # Orange left accent bar
    add_shape_bg(slide, left, top, Inches(0.08), height, ORANGE)
    # Text
    add_textbox(slide, left + Inches(0.25), top + Inches(0.1),
                width - Inches(0.35), height - Inches(0.2),
                text, font_size=font_size, bold=True, color=DARK_BLUE,
                alignment=PP_ALIGN.LEFT)
    return bg


# ---------------------------------------------------------------------------
# SLIDE BUILDERS — v6 (26 slides)
# ---------------------------------------------------------------------------

def slide_01_title(prs):
    """SLIDE 1: Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BLUE)

    # Orange accent bar at top
    add_shape_bg(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.15), ORANGE)

    # Main title
    add_textbox(slide, Inches(1.5), Inches(1.6), Inches(10.3), Inches(1.2),
                "QPH Multi-Idioma:", font_size=Pt(44), bold=True,
                color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.5), Inches(2.7), Inches(10.3), Inches(1.2),
                "Problematicas, Datos y Plan de Accion", font_size=Pt(36), bold=True,
                color=ORANGE, alignment=PP_ALIGN.CENTER)

    # Subtitle
    add_textbox(slide, Inches(1.5), Inches(4.1), Inches(10.3), Inches(0.6),
                "Reunion de Equipo \u00b7 20 Feb 2026", font_size=Pt(22),
                bold=False, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Confidential
    add_textbox(slide, Inches(1.5), Inches(5.1), Inches(10.3), Inches(0.5),
                "Confidencial \u2014 Solo uso interno", font_size=Pt(14),
                bold=False, color=RGBColor(0x99, 0xAA, 0xCC),
                alignment=PP_ALIGN.CENTER)

    # Bottom orange accent bar
    add_shape_bg(slide, Inches(0), Inches(7.35), SLIDE_WIDTH, Inches(0.15), ORANGE)

    add_slide_number(slide, 1)


def slide_02_agenda(prs):
    """SLIDE 2: Agenda."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_GRAY)
    add_slide_header(slide, "Agenda")

    items = [
        "1.  Donde estamos (datos reales)",
        "2.  Los 6 problemas criticos",
        "3.  El caso del 'Cerdo Robot' (ejemplo real)",
        "4.  La solucion propuesta",
        "5.  Objetivos SMART por persona",
        "6.  4 Investigaciones Profundas",
        "7.  Siguiente paso",
    ]
    add_bullet_list(slide, Inches(1.5), Inches(1.8), Inches(10.0), Inches(5.2),
                    items, font_size=Pt(24), color=DARK_BLUE,
                    bold_items={0, 1, 2, 3, 4, 5, 6})

    add_slide_number(slide, 2)


def slide_03_panorama(prs):
    """SLIDE 3: El Panorama — Datos Reales del Canal."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_GRAY)
    add_slide_header(slide, "El Panorama \u2014 Datos Reales del Canal",
                     "Donde estamos hoy")

    items = [
        "27 idiomas doblados con ElevenLabs",
        "Solo 1 idioma (EN) tiene revision humana — 26 se publican a ciegas",
        "ES = ~58% del ingreso  |  EN = ~22%  |  PT = ~6%  |  Top 3 = ~86% del ingreso",
        "Los otros 21 idiomas juntos = ~7% del ingreso",
        "1 view desde Alemania vale x7.2 mas que 1 view desde Mexico",
        "1 view desde EEUU vale x2.4 mas que 1 view desde Mexico",
        "El canal AUN no alcanza break-even mensual",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(4.2),
                    items, font_size=Pt(19), color=DARK_GRAY,
                    bold_items={2, 3, 4, 5, 6})

    add_highlight_box(slide, Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.85),
                      "Top 3 idiomas (ES+EN+PT) = 86% del ingreso. "
                      "Los otros 21 idiomas combinados = 7%",
                      font_size=Pt(18))

    add_slide_number(slide, 3)


def slide_04_tabla_idiomas(prs):
    """SLIDE 4: Ingreso Estimado por Idioma (tabla principal)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_GRAY)
    add_slide_header(slide, "Ingreso Estimado por Idioma",
                     "Ordenado por importancia de revenue — mayor a menor")

    rows = [
        ["Idioma", "Views", "% Ingreso est.", "Valor/view*", "AVD", "Delta vs ES", "Prioridad"],
        ["ES (original)", "232.2M", "~58%", "x1.0 (base)", "4:41", "baseline", "BASE — proteger"],
        ["EN (master)", "40.8M", "~22%", "x2.4", "3:04", "-35%", "CRITICA — afecta 26 idiomas"],
        ["PT", "35.5M", "~6%", "x0.7", "3:49", "-18%", "ALTA — 2do mercado x volumen"],
        ["DE", "2.2M", "~4%", "x7.2", "3:40", "-22%", "PREMIUM — crecer volumen"],
        ["IT", "4.4M", "~3%", "x2.9", "3:34", "-24%", "ALTA — buen RPM + volumen"],
        ["FR", "4.9M", "~2%", "x1.8", "3:38", "-22%", "MEDIA — Europa + Africa FR"],
        ["RU", "19.0M", "~1%", "x0.15", "3:43", "-21%", "BAJA — VPN? revenue no cuadra"],
        ["TR", "5.4M", "~0.7%", "x0.5", "3:33", "-24%", "BAJA — emergente"],
        ["JA", "692K", "~0.4%", "x2.1", "2:22", "-49%", "POTENCIAL — premium, AVD muy bajo"],
        ["KO", "629K", "~0.2%", "x1.2", "2:27", "-48%", "EVALUAR — premium, poco vol"],
        ["ID", "3.6M", "~0.3%", "x0.3", "3:06", "-34%", "EVALUAR — mucho vol, poco RPM"],
        ["HI+TA+FIL+ZH+MS", "2.0M", "~0.3%", "x0.3-0.4", "~2:20", "-50%", "PAUSAR/EVALUAR"],
    ]
    build_table(
        slide, Inches(0.25), Inches(1.45), Inches(12.83), Inches(5.1), rows,
        col_widths=[Inches(1.8), Inches(1.1), Inches(1.2), Inches(1.1),
                    Inches(0.85), Inches(1.05), Inches(5.72)],
        font_size=Pt(10))

    add_textbox(slide, Inches(0.25), Inches(6.6), Inches(12.83), Inches(0.28),
                "*Valor/view = RPM estimado del idioma / RPM del espanol. "
                "Metodologia: pais como proxy de idioma + cruce con audio tracks.",
                font_size=Pt(8.5), bold=False, color=RGBColor(0x88, 0x88, 0x88))

    add_highlight_box(slide, Inches(0.25), Inches(6.9), Inches(12.83), Inches(0.5),
                      "DE vale x7.2 por view pero tiene solo 2.2M views. "
                      "Si doblamos bien el aleman y crecemos volumen, cada view nueva vale 7x mas.",
                      font_size=Pt(13))

    add_slide_number(slide, 4)


def slide_05_perdida_audiencia(prs):
    """SLIDE 5: Lo que Estamos Perdiendo — Perdida de Audiencia (AVD bars)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_GRAY)
    add_slide_header(slide, "Lo que Estamos Perdiendo \u2014 Perdida de Audiencia",
                     "AVD = Duracion Promedio de Vista. ES original = 4:41 = 100%")

    # Bar chart data: (label, avd_str, pct_of_es, color)
    # pct_of_es: percentage of ES AVD (4:41 = 281 seconds baseline)
    bars = [
        ("ES (original)",   "4:41", 100, GREEN),
        ("PT-BR",           "3:49",  81, GREEN),
        ("Europeos (DE,FR)","3:37",  77, GREEN),
        ("IT",              "3:34",  76, GREEN),
        ("RU / TR",         "3:38",  77, GREEN),
        ("EN (master)",     "3:04",  65, ORANGE),
        ("Arabe",           "3:23",  72, ORANGE),
        ("Hindi",           "2:56",  63, ORANGE),
        ("CJK (JA,KO,ZH)", "2:24",  51, RED),
        ("SEA",             "2:39",  57, RED),
        ("Tamil",           "1:58",  42, RED),
    ]

    label_col_w = Inches(2.0)
    avd_col_w   = Inches(0.6)
    bar_max_w   = Inches(7.8)
    pct_col_w   = Inches(0.7)
    row_h       = Inches(0.42)
    start_x     = Inches(0.4)
    bar_start_x = start_x + label_col_w + avd_col_w + Inches(0.1)
    start_y     = Inches(1.5)

    for i, (label, avd, pct, color) in enumerate(bars):
        y = start_y + i * row_h

        # Label
        add_textbox(slide, start_x, y + Inches(0.04), label_col_w, row_h - Inches(0.05),
                    label, font_size=Pt(11), bold=(i == 0 or i == 5),
                    color=DARK_GRAY if i != 0 else DARK_BLUE)

        # AVD value
        add_textbox(slide, start_x + label_col_w, y + Inches(0.04),
                    avd_col_w, row_h - Inches(0.05),
                    avd, font_size=Pt(11), bold=True, color=color,
                    alignment=PP_ALIGN.CENTER)

        # Bar
        bar_w = bar_max_w * pct / 100.0
        add_shape_bg(slide, bar_start_x, y + Inches(0.06),
                     bar_w, row_h - Inches(0.12), color)

        # Percentage label after bar
        add_textbox(slide, bar_start_x + bar_w + Inches(0.05), y + Inches(0.04),
                    pct_col_w, row_h - Inches(0.05),
                    f"{pct}%", font_size=Pt(10), bold=False, color=DARK_GRAY)

    # Legend
    legend_y = Inches(6.25)
    for leg_color, leg_label in [(GREEN, ">75%: saludable"),
                                  (ORANGE, "60-75%: alerta"),
                                  (RED, "<60%: critico")]:
        pass  # drawn below
    add_shape_bg(slide, Inches(0.4), legend_y, Inches(0.25), Inches(0.2), GREEN)
    add_textbox(slide, Inches(0.7), legend_y - Inches(0.02), Inches(2.0), Inches(0.25),
                ">75%: saludable", font_size=Pt(10), color=DARK_GRAY)
    add_shape_bg(slide, Inches(3.0), legend_y, Inches(0.25), Inches(0.2), ORANGE)
    add_textbox(slide, Inches(3.3), legend_y - Inches(0.02), Inches(2.0), Inches(0.25),
                "60-75%: alerta", font_size=Pt(10), color=DARK_GRAY)
    add_shape_bg(slide, Inches(5.6), legend_y, Inches(0.25), Inches(0.2), RED)
    add_textbox(slide, Inches(5.9), legend_y - Inches(0.02), Inches(2.0), Inches(0.25),
                "<60%: critico", font_size=Pt(10), color=DARK_GRAY)

    add_highlight_box(slide, Inches(0.4), Inches(6.55), Inches(12.53), Inches(0.7),
                      "EN pierde 35% y es el master de 26 idiomas. CJK pierde 49%. "
                      "Si mejoramos AVD de EN de 3:04 a 4:00, el impacto se multiplica x26.",
                      font_size=Pt(15))

    add_slide_number(slide, 5)


def slide_06_problema1(prs):
    """SLIDE 6: Problema 1 — Telefono Descompuesto."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_GRAY)
    add_slide_header(slide, "Problema 1 \u2014 Telefono Descompuesto",
                     "La cadena de traduccion multiplica cada error x26")

    # Chain diagram: ES -> EN -> 26 idiomas
    box_h = Inches(0.9)
    box_y = Inches(1.9)

    es_box = add_shape_bg(slide, Inches(0.8), box_y, Inches(2.2), box_h, GREEN)
    set_text(es_box.text_frame, "ES (Original)", font_size=Pt(18), bold=True,
             color=WHITE, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(3.05), box_y, Inches(0.7), box_h,
                "\u2192", font_size=Pt(36), bold=True, color=ORANGE,
                alignment=PP_ALIGN.CENTER)

    en_box = add_shape_bg(slide, Inches(3.75), box_y, Inches(2.5), box_h, ORANGE)
    set_text(en_box.text_frame, "EN (Master)", font_size=Pt(18), bold=True,
             color=WHITE, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(6.3), box_y, Inches(0.7), box_h,
                "\u2192", font_size=Pt(36), bold=True, color=RED,
                alignment=PP_ALIGN.CENTER)

    lang_box = add_shape_bg(slide, Inches(7.0), box_y, Inches(3.2), box_h, RED)
    set_text(lang_box.text_frame, "26 idiomas", font_size=Pt(18), bold=True,
             color=WHITE, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(7.0), Inches(2.85), Inches(3.2), Inches(0.45),
                "Cada error x26", font_size=Pt(13), bold=True, color=RED,
                alignment=PP_ALIGN.CENTER)

    # Example
    add_textbox(slide, Inches(0.8), Inches(3.5), Inches(11.5), Inches(0.5),
                'Ejemplo real: "Que padre!" \u2192 EN: "What a father!" \u2192 propagado a 26 idiomas',
                font_size=Pt(16), bold=False, color=DARK_GRAY)

    items = [
        "EN master ya pierde 35% de AVD (3:04 vs ES 4:41) — parte del problema",
        "Cada error de traduccion se multiplica x26 idiomas",
        "Expresiones coloquiales mexicanas se traducen literalmente sin sentido",
        "5 de 5 opiniones AI coinciden: este es el problema #1",
        "Hipotesis: PT tiene AVD de 3:49 (-18%) — podria mejorar con ruta directa ES\u2192PT",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(4.15), Inches(11.5), Inches(2.4),
                    items, font_size=Pt(17), color=DARK_GRAY,
                    bold_items={0, 1, 3})

    add_slide_number(slide, 6)


def slide_07_problema2(prs):
    """SLIDE 7: Problema 2 — Zero Metricas de Calidad."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_GRAY)
    add_slide_header(slide, "Problema 2 \u2014 Zero Metricas de Calidad",
                     "No se mide = no se puede mejorar")

    # Broken gauge
    gauge_shape = slide.shapes.add_shape(
        MSO_SHAPE.DONUT, Inches(10.2), Inches(1.9), Inches(2.0), Inches(2.0))
    gauge_shape.fill.solid()
    gauge_shape.fill.fore_color.rgb = RED
    gauge_shape.line.color.rgb = DARK_GRAY
    gauge_shape.line.width = Pt(2)

    add_textbox(slide, Inches(10.55), Inches(2.2), Inches(1.3), Inches(1.3),
                "X", font_size=Pt(48), bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    rows = [
        ["Metrica", "Estado actual"],
        ["WER (tasa de error de traduccion)", "No la medimos"],
        ["Tasa de error por idioma", "No la tenemos"],
        ["Tiempo de dubbing completo", "No se sabe"],
        ["Tasa de aprobacion a la primera (FTR)", "Desconocida"],
        ["Idiomas con revision humana", "Solo 1 de 27 (EN)"],
        ["Idiomas que se publican sin verificacion", "26 de 27"],
    ]
    build_table(
        slide, Inches(0.5), Inches(1.5), Inches(9.5), Inches(4.2), rows,
        col_widths=[Inches(5.5), Inches(4.0)],
        font_size=Pt(13))

    add_quote_box(slide, Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.75),
                  '"5 de 5 opiniones AI coinciden: sin medir, no se puede mejorar. '
                  'Actualmente volamos a ciegas en 26 idiomas."',
                  font_size=Pt(17))

    add_highlight_box(slide, Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.55),
                      "26 idiomas se publican sin ninguna verificacion de calidad",
                      font_size=Pt(17))

    add_slide_number(slide, 7)


def slide_08_problema3(prs):
    """SLIDE 8: Problema 3 — Blacklists: Solo 13 Palabras para 27 Idiomas."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_GRAY)
    add_slide_header(slide, "Problema 3 \u2014 Blacklists: Solo 13 Terminos para 27 Idiomas",
                     "Contenido para ninos 8-15 anos sin filtro automatico")

    # Real data from blacklist JSON files
    rows = [
        ["Archivo", "Terminos", "Palabras bloqueadas", "Estado en pipeline"],
        ["blacklist_global.json", "6 palabras",
         "muerte, matar, sexy, sangre, droga, suicidio",
         "NO integrado — referencia manual"],
        ["blacklist_ar.json (Arabe)", "5 terminos",
         "cerdo/khanzir, alcohol, dios, iglesia, beso",
         "NO integrado — referencia manual"],
        ["blacklist_de.json (Aleman)", "2 terminos",
         "nazi, guerra",
         "NO integrado — referencia manual"],
        ["24 idiomas restantes", "0 terminos",
         "Sin blacklist de ningun tipo",
         "VACIO — cero cobertura"],
    ]
    build_table(
        slide, Inches(0.3), Inches(1.5), Inches(12.73), Inches(2.9), rows,
        col_widths=[Inches(2.5), Inches(1.4), Inches(4.73), Inches(4.1)],
        font_size=Pt(11))

    # Visual: 3 green + 24 red boxes
    box_size = Inches(0.33)
    start_x = Inches(0.7)
    box_y = Inches(4.6)
    for i in range(27):
        x = start_x + i * Inches(0.45)
        c = GREEN if i < 3 else RED
        add_shape_bg(slide, x, box_y, box_size, box_size, c)

    add_textbox(slide, Inches(0.7), Inches(5.0), Inches(12.0), Inches(0.35),
                "3 idiomas con blacklist parcial (sin integrar)                    "
                "24 idiomas con CERO blacklist",
                font_size=Pt(11), bold=False, color=DARK_GRAY)

    items = [
        "Los archivos existen pero NO estan integrados al pipeline — todo es manual o inexistente",
        "El directorio by_language/ referenciado en el JSON global NO EXISTE",
    ]
    add_bullet_list(slide, Inches(0.3), Inches(5.45), Inches(12.73), Inches(0.9),
                    items, font_size=Pt(14), color=DARK_GRAY, bold_items={0, 1})

    add_highlight_box(slide, Inches(0.3), Inches(6.5), Inches(12.73), Inches(0.75),
                      "TOTAL: 13 terminos en 3 archivos. 24 idiomas con CERO blacklist. "
                      "Ninguno esta conectado al pipeline automatico.",
                      font_size=Pt(15))

    add_slide_number(slide, 8)


def slide_09_problema4(prs):
    """SLIDE 9: Problema 4 — Desconexion Cultural."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_GRAY)
    add_slide_header(slide, "Problema 4 \u2014 Desconexion Cultural",
                     "No existe un diccionario cultural positivo")

    rows = [
        ["Cluster", "Problema", "Ejemplo", "Riesgo"],
        ["Arabe", "Animal haram como insulto", '"cerdo robot" = ofensa grave', "CRITICO"],
        ["CJK (JA/KO/ZH)", "Formalidad incorrecta", "Ninos hablando como ejecutivos", "ALTO"],
        ["Japones", "4,500 onomatopeyas", '"guau" no es "wan wan"', "ALTO"],
        ["Aleman", "du vs Sie", 'Ninos de "usted" entre si', "MEDIO"],
        ["Frances", "tu vs vous", "Amigos formales entre si", "MEDIO"],
        ["Hindi", "tum vs aap", "Registro adulto en contenido infantil", "MEDIO"],
        ["Global", "Humor mexicano", '"No manches" = "Don\'t stain"', "ALTO"],
    ]
    build_table(
        slide, Inches(0.3), Inches(1.45), Inches(12.73), Inches(3.0), rows,
        col_widths=[Inches(2.0), Inches(3.3), Inches(4.43), Inches(3.0)],
        font_size=Pt(11))

    # Expression adaptation sub-table
    add_textbox(slide, Inches(0.3), Inches(4.55), Inches(12.73), Inches(0.35),
                "Ejemplos de adaptacion de expresiones mexicanas:", font_size=Pt(13),
                bold=True, color=DARK_BLUE)

    rows2 = [
        ["Expresion MX", "Significado", "EN", "JA", "AR", "DE"],
        ['"Que padre!"', "Que genial!", "That's awesome!", "Sugoi!", "Yaa salaam!", "Mega cool!"],
        ['"No manches"', "No puede ser", "No way!", "Maji de?!", "Mish ma'qul!", "Krass!"],
    ]
    build_table(
        slide, Inches(0.3), Inches(4.9), Inches(12.73), Inches(1.3), rows2,
        col_widths=[Inches(1.7), Inches(1.6), Inches(2.0), Inches(2.0), Inches(2.73), Inches(2.7)],
        font_size=Pt(11))

    add_highlight_box(slide, Inches(0.3), Inches(6.3), Inches(12.73), Inches(0.6),
                      "No existe diccionario cultural positivo. Solo tenemos lo que NO decir "
                      "(13 palabras) pero no lo que SI decir.",
                      font_size=Pt(15))

    add_slide_number(slide, 9)


def slide_10_problema5(prs):
    """SLIDE 10: Problema 5 — Dos Sistemas Desconectados."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_GRAY)
    add_slide_header(slide, "Problema 5 \u2014 Dos Sistemas Desconectados",
                     "El motor y las ruedas no tienen conexion")

    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5),
                '"Tenemos dos mitades que no se hablan"',
                font_size=Pt(20), bold=True, color=DARK_BLUE)

    # Left box: API Layer
    api_box = add_shape_bg(slide, Inches(0.8), Inches(2.2), Inches(4.5),
                           Inches(2.5), DARK_BLUE)
    set_text(api_box.text_frame, "API Layer", font_size=Pt(22), bold=True,
             color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.0), Inches(3.05), Inches(4.1), Inches(1.5),
                "Sube videos\nGestiona ElevenLabs\nOrquesta produccion",
                font_size=Pt(16), color=WHITE)

    # Disconnect
    add_textbox(slide, Inches(5.4), Inches(2.85), Inches(2.4), Inches(1.5),
                "\u26a0\nDESCONECTADOS", font_size=Pt(18), bold=True,
                color=RED, alignment=PP_ALIGN.CENTER)

    # Right box: ERP Layer
    erp_box = add_shape_bg(slide, Inches(8.0), Inches(2.2), Inches(4.5),
                           Inches(2.5), LIGHT_BLUE)
    set_text(erp_box.text_frame, "ERP Layer", font_size=Pt(22), bold=True,
             color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(8.2), Inches(3.05), Inches(4.1), Inches(1.5),
                "Calcula errores\nEstima costos\nReportes internos",
                font_size=Pt(16), color=WHITE)

    # Bugs
    add_textbox(slide, Inches(0.8), Inches(4.85), Inches(11.7), Inches(0.35),
                "Bugs conocidos:", font_size=Pt(14), bold=True, color=DARK_GRAY)
    items = [
        "Bug P0 (critico): prescanner crashea cuando recibe dato inesperado",
        "Bug P1 (alto): toda medicion de error se guarda como 'espanol' sin importar el idioma real",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.0),
                    items, font_size=Pt(15), color=DARK_GRAY, bold_items={0, 1})

    add_quote_box(slide, Inches(0.8), Inches(6.35), Inches(11.7), Inches(0.75),
                  '"Es como tener el motor y las ruedas sin conexion — '
                  'los datos existen pero no se comunican entre sistemas."',
                  font_size=Pt(16))

    add_slide_number(slide, 10)


def slide_11_problema6(prs):
    """SLIDE 11: Problema 6 — Inversion a Ciegas (Sin ROI por Idioma)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_GRAY)
    add_slide_header(slide, "Problema 6 \u2014 Inversion a Ciegas",
                     "Nunca se ha calculado ROI por idioma")

    items = [
        "Nunca se ha calculado cuanto cuesta producir cada idioma vs cuanto retorna",
        "El canal AUN no esta en break-even — cada episodio tiene costo sin retorno garantizado",
        "Concentracion: ES+EN+PT = ~86% del ingreso | Los otros 21 idiomas = ~7%",
        "Cada idioma tiene costo de dubbing en ElevenLabs + tiempo de revision humana",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(2.5),
                    items, font_size=Pt(18), color=DARK_GRAY, bold_items={2, 3})

    add_textbox(slide, Inches(0.8), Inches(4.1), Inches(11.7), Inches(0.4),
                "Candidatos a pausar (ROI negativo estimado):",
                font_size=Pt(15), bold=True, color=DARK_BLUE)

    rows = [
        ["Idioma", "Views", "AVD", "Situacion", "Decision pendiente"],
        ["Tamil", "30K", "1:58", "AVD 42% del ES — peor caso", "Pausar?"],
        ["Filipino", "46K", "~2:30", "Volumen y AVD bajos", "Pausar?"],
        ["Chino (ZH)", "65K", "~2:20", "Revenue ~cero (VPN/restricciones)", "Pausar?"],
        ["Hindi", "~800K", "2:56", "Volumen ok, AVD bajo, RPM muy bajo", "Degradar a Tier 3?"],
    ]
    build_table(
        slide, Inches(0.8), Inches(4.55), Inches(11.7), Inches(2.0), rows,
        col_widths=[Inches(1.8), Inches(1.2), Inches(1.0), Inches(4.7), Inches(3.0)],
        font_size=Pt(12))

    add_highlight_box(slide, Inches(0.8), Inches(6.65), Inches(11.7), Inches(0.65),
                      "Sin calcular ROI por idioma, no sabemos si estamos ganando o perdiendo "
                      "en 21 de los 27 idiomas.",
                      font_size=Pt(16))

    add_slide_number(slide, 11)


def slide_12_cerdo_robot(prs):
    """SLIDE 12: El Caso del 'Cerdo Robot' — Ejemplo Real."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_GRAY)
    add_slide_header(slide, "El Caso del 'Cerdo Robot' \u2014 Ejemplo Real",
                     "EP042 — El quasi-incidente que casi pasa a produccion")

    # Step-by-step flow
    steps = [
        ("1. Guion ES", DARK_BLUE,
         '"No manches, este lab esta bien canon!\nMiren, hasta tiene un cerdo robot!"'),
        ("2. ElevenLabs (sin corregir)", RED,
         '"Don\'t stain, this lab is very cannon!\nLook, it even has a pig robot!" \u2014 literal, sin sentido'),
        ("3. Saul/Ivan corrigen EN", GREEN,
         '"No way, this lab is so cool!\nLook, it even has a robot pig!"'),
        ("4. Del EN se genera AR", RED,
         '"robot pig" \u2192 "khanzir ali" \u2014 khanzir = CERDO en arabe\n= OFENSA CULTURAL GRAVE para audiencia musulmana'),
        ("5. Nadie lo detecta", RED,
         "El arabe no se revisa. Se publica. Llegan comentarios negativos."),
    ]

    box_w = Inches(2.35)
    box_h = Inches(1.55)
    box_y = Inches(1.45)
    gap = Inches(0.15)

    for i, (step_title, color, step_text) in enumerate(steps):
        x = Inches(0.2) + i * (box_w + gap)

        # Step box
        step_box = add_shape_bg(slide, x, box_y, box_w, box_h, color)
        tf = step_box.text_frame
        tf.word_wrap = True
        set_text(tf, step_title, font_size=Pt(12), bold=True,
                 color=WHITE, alignment=PP_ALIGN.CENTER)
        add_paragraph(tf, step_text, font_size=Pt(9), bold=False,
                      color=WHITE, alignment=PP_ALIGN.LEFT,
                      space_before=Pt(4))

        # Arrow between steps
        if i < 4:
            arrow_x = x + box_w + Inches(0.01)
            add_textbox(slide, arrow_x, box_y + Inches(0.55), gap + Inches(0.02),
                        Inches(0.5), "\u2192", font_size=Pt(20), bold=True,
                        color=ORANGE, alignment=PP_ALIGN.CENTER)

    # Error table
    add_textbox(slide, Inches(0.3), Inches(3.15), Inches(12.73), Inches(0.35),
                "Errores detectados post-publicacion:", font_size=Pt(13),
                bold=True, color=DARK_BLUE)

    rows = [
        ["Idioma", "Error detectado", "Categoria"],
        ["Arabe", '"khanzir" (cerdo) = ofensa cultural grave para audiencia musulmana', "A — CRITICO"],
        ["Japones", '"No manches" sin equivalente cultural — suena extrano', "B"],
        ["Coreano", "Gabriel habla con formalidad corporativa (no infantil)", "C"],
        ["Aleman", "Timing desbordado 2.3 segundos — audio corta", "Info"],
    ]
    build_table(
        slide, Inches(0.3), Inches(3.5), Inches(12.73), Inches(2.4), rows,
        col_widths=[Inches(1.5), Inches(8.73), Inches(2.5)],
        font_size=Pt(12))

    add_highlight_box(slide, Inches(0.3), Inches(6.0), Inches(12.73), Inches(0.85),
                      "Un solo modismo mexicano ('cerdo robot') se convierte en ofensa cultural en arabe. "
                      "Sin blacklists integradas, nadie lo detecta antes de publicar.",
                      font_size=Pt(16))

    add_slide_number(slide, 12)


def slide_13_buena_noticia(prs):
    """SLIDE 13: La Buena Noticia — 60% ya construido."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_GRAY)
    add_slide_header(slide, "La Buena Noticia",
                     "Ya tenemos avance significativo — 60% del codigo construido")

    # Big percentage
    add_textbox(slide, Inches(4.5), Inches(1.6), Inches(4.3), Inches(1.5),
                "60%", font_size=Pt(72), bold=True, color=GREEN,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(3.5), Inches(3.0), Inches(6.3), Inches(0.6),
                "del codigo ya esta construido",
                font_size=Pt(22), bold=False, color=DARK_GRAY,
                alignment=PP_ALIGN.CENTER)

    # Phase status boxes
    phases = [
        ("Phase 1: Pre-produccion", "LISTO", GREEN),
        ("Phase 2: Produccion", "LISTO", GREEN),
        ("Phase 2.5: ERP", "PARCIAL", ORANGE),
        ("Phase 3: QA Automatizado", "PENDIENTE", RED),
    ]

    for i, (label, status, color) in enumerate(phases):
        y = Inches(3.9) + i * Inches(0.65)
        add_textbox(slide, Inches(2.0), y, Inches(5.5), Inches(0.55),
                    label, font_size=Pt(18), bold=False, color=DARK_GRAY)
        badge = add_shape_bg(slide, Inches(8.0), y + Inches(0.05),
                             Inches(2.5), Inches(0.45), color)
        set_text(badge.text_frame, status, font_size=Pt(14), bold=True,
                 color=WHITE, alignment=PP_ALIGN.CENTER)

    add_quote_box(slide, Inches(2.0), Inches(6.7), Inches(9.3), Inches(0.6),
                  '"Solo falta Phase 3: QA Automatizado + conectar todo"',
                  font_size=Pt(17))

    add_slide_number(slide, 13)


def slide_14_solucion(prs):
    """SLIDE 14: La Solucion — 4 Gates de Calidad."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_GRAY)
    add_slide_header(slide, "La Solucion \u2014 4 Gates de Calidad",
                     "Cada episodio pasa por 4 puntos de control antes de publicarse")

    gates = [
        ("Gate 1", "Pre-scan del guion",
         "ANTES de ElevenLabs\nDetecta blacklists\n+ modismos", GREEN),
        ("Gate 2", "Verificacion post-dubbing",
         "WER + timing\npor idioma", ORANGE),
        ("Gate 3", "Auditoria multi-IA",
         "Claude + GPT-4\n+ Gemini votan", LIGHT_BLUE),
        ("Gate 4", "Dashboard humano",
         "Decision final\ndel equipo", DARK_BLUE),
    ]

    gate_width = Inches(2.6)
    gap = Inches(0.3)
    start_x = Inches(0.7)

    for i, (gate_id, title, desc, color) in enumerate(gates):
        x = start_x + i * (gate_width + gap)
        y_top = Inches(1.75)

        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Inches(0.8), y_top, Inches(1.0), Inches(1.0))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        set_text(circle.text_frame, gate_id, font_size=Pt(15), bold=True,
                 color=WHITE, alignment=PP_ALIGN.CENTER)

        box = add_shape_bg(slide, x, y_top + Inches(1.2), gate_width,
                           Inches(2.1), color)
        tf = box.text_frame
        tf.word_wrap = True
        set_text(tf, title, font_size=Pt(16), bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
        add_paragraph(tf, desc, font_size=Pt(13), bold=False, color=WHITE,
                      alignment=PP_ALIGN.CENTER, space_before=Pt(8))

        if i < 3:
            arrow_x = x + gate_width + Inches(0.02)
            add_textbox(slide, arrow_x, y_top + Inches(1.9), gap, Inches(0.5),
                        "\u2192", font_size=Pt(28), bold=True, color=ORANGE,
                        alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.8),
                "Flujo: Guion \u2192 Gate 1 (pre-scan ANTES de ElevenLabs) \u2192 "
                "ElevenLabs dubbing \u2192 Gate 2 (WER+timing) \u2192 "
                "Gate 3 (multi-IA) \u2192 Gate 4 (humano aprueba) \u2192 Publicar",
                font_size=Pt(15), bold=False, color=DARK_GRAY)

    add_highlight_box(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.6),
                      "Gate 1 detiene errores culturales ANTES de gastar en dubbing. "
                      "Gate 4 da decision final al equipo humano.",
                      font_size=Pt(15))

    add_slide_number(slide, 14)


def slide_15_costo(prs):
    """SLIDE 15: Costo del Sistema."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_GRAY)
    add_slide_header(slide, "Costo del Sistema",
                     "Todo relativo — sin cifras absolutas")

    rows = [
        ["Escenario", "Cobertura", "Costo relativo", "Resultado"],
        ["Hoy (sin sistema)",
         "Solo EN revisado\n(1 de 27)",
         "~100% del costo va\na 1 idioma",
         "26 idiomas sin control"],
        ["Revision humana completa\n(sin automatizar)",
         "27 idiomas con\nrevision humana",
         "~100x el costo actual\n(insostenible)",
         "Imposible a escala"],
        ["Solucion propuesta\n(automatico + Tier 1 humano)",
         "27 idiomas cubiertos\n(5 con revision humana)",
         "~1.2x el costo actual\n(centavos de API)",
         "27 idiomas con control"],
    ]
    build_table(
        slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(3.0), rows,
        col_widths=[Inches(3.0), Inches(2.7), Inches(2.8), Inches(3.8)],
        font_size=Pt(13))

    items = [
        "El costo automatico por episodio son centavos de dolar en API de IA",
        "Con el sistema: pasamos de revisar 1 idioma a revisar los 27",
        "Revision humana se concentra en Tier 1 (5 idiomas = ~37% del ingreso no-ES)",
        "Tier 2 y Tier 3 se cubren con muestreo automatico — casi cero costo marginal",
    ]
    add_bullet_list(slide, Inches(0.5), Inches(4.65), Inches(12.3), Inches(2.0),
                    items, font_size=Pt(17), color=DARK_GRAY, bold_items={1, 2})

    add_highlight_box(slide, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.65),
                      "Por ~1.2x el costo actual cubrimos 27 idiomas en vez de 1.",
                      font_size=Pt(18))

    add_slide_number(slide, 15)


def slide_16_tiers(prs):
    """SLIDE 16: Prioridad por Tiers."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_GRAY)
    add_slide_header(slide, "Prioridad por Tiers",
                     "Invertir donde mas retorna, automatizar el resto")

    tiers = [
        ("Tier 1", "5 idiomas", "EN, PT, DE, FR, IT",
         "Revision humana completa — ~37% del ingreso no-ES", GREEN),
        ("Tier 2", "5 idiomas", "RU, TR, AR, ID, JA",
         "Muestreo automatico 30% de episodios", ORANGE),
        ("Tier 3", "17 idiomas", "Resto (KO, HI, ZH, TA, FIL, MS...)",
         "Solo automatico — alertas si algo supera umbral", LIGHT_BLUE),
    ]

    for i, (tier, count, langs, method, color) in enumerate(tiers):
        y = Inches(1.75) + i * Inches(1.65)

        badge = add_shape_bg(slide, Inches(0.5), y, Inches(2.2),
                             Inches(1.3), color)
        tf = badge.text_frame
        tf.word_wrap = True
        set_text(tf, tier, font_size=Pt(22), bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
        add_paragraph(tf, count, font_size=Pt(14), color=WHITE,
                      alignment=PP_ALIGN.CENTER)

        add_textbox(slide, Inches(3.1), y + Inches(0.1), Inches(9.7),
                    Inches(0.5), langs, font_size=Pt(18), bold=True,
                    color=DARK_GRAY)
        add_textbox(slide, Inches(3.1), y + Inches(0.65), Inches(9.7),
                    Inches(0.5), method, font_size=Pt(15), bold=False,
                    color=DARK_GRAY)

    add_quote_box(slide, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.7),
                  '"Invertir donde mas retorna, automatizar el resto. '
                  'Tier 1 = ~37% del ingreso no-ES con revision humana completa."',
                  font_size=Pt(16))

    add_slide_number(slide, 16)


def slide_17_smart_overview(prs):
    """SLIDE 17: Objetivos SMART del Equipo — Vista General."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_GRAY)
    add_slide_header(slide, "Objetivos SMART del Equipo \u2014 Vista General",
                     "Especificos, Medibles, Alcanzables, Relevantes, con Tiempo")

    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.5),
                "Cada persona tiene responsabilidades claras y metas medibles a 90 dias:",
                font_size=Pt(18), bold=False, color=DARK_GRAY)

    people = [
        ("Andrea", "Contenido y Adaptacion Cultural",
         "Blacklists, expresiones, revision de guion pre-dubbing", GREEN),
        ("Alan", "Optimizacion Post-Produccion",
         "Eficiencia de Saul/Ivan, identificacion de gaps, coordinacion", ORANGE),
        ("Ramon", "Soporte Tecnico del Pipeline",
         "Herramientas, automatizacion, infraestructura tecnica", LIGHT_BLUE),
        ("Daniel", "Deep Research y Habilitacion IA",
         "4 investigaciones profundas, herramientas AI, apoyo al equipo", DARK_BLUE),
    ]

    card_width = Inches(5.7)
    card_height = Inches(1.05)
    col1_x = Inches(0.8)
    col2_x = Inches(6.8)

    for i, (name, role, desc, color) in enumerate(people):
        x = col1_x if i % 2 == 0 else col2_x
        y = Inches(2.2) + (i // 2) * Inches(1.3)

        card = add_shape_bg(slide, x, y, card_width, card_height, color)
        tf = card.text_frame
        tf.word_wrap = True
        set_text(tf, f"{name} \u2014 {role}", font_size=Pt(15), bold=True,
                 color=WHITE, alignment=PP_ALIGN.LEFT)
        add_paragraph(tf, desc, font_size=Pt(12), bold=False, color=WHITE,
                      alignment=PP_ALIGN.LEFT, space_before=Pt(4))

    add_quote_box(slide, Inches(0.8), Inches(4.9), Inches(11.7), Inches(0.65),
                  '"Los primeros 30 dias son para MEDIR. '
                  'No podemos mejorar lo que no medimos."',
                  font_size=Pt(17))

    rows = [
        ["Metrica", "Hoy", "Meta 30 dias", "Meta 90 dias"],
        ["Idiomas con metricas", "1 de 27", "5 (Tier 1)", "27 de 27"],
        ["Blacklists integradas al pipeline", "0 de 27", "5 (Tier 1)", "27 de 27"],
        ["Episodios aprobados 1ra vez (FTR)", "Desconocido", "Medido", "> 60%"],
        ["Incidentes Cat. A en Tier 1", "No medido", "0 meta", "0 siempre"],
    ]
    build_table(slide, Inches(0.8), Inches(5.65), Inches(11.7), Inches(1.6), rows,
                col_widths=[Inches(4.0), Inches(2.0), Inches(2.8), Inches(2.9)],
                font_size=Pt(10))

    add_slide_number(slide, 17)


def slide_18_smart_andrea(prs):
    """SLIDE 18: SMART — Andrea (Contenido y Adaptacion Cultural)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_GRAY)
    add_slide_header(slide, "Andrea \u2014 Contenido y Adaptacion Cultural",
                     "Guionismo, blacklists, temas culturales")

    rows = [
        ["#", "Objetivo", "Como se mide", "Meta", "Plazo"],
        ["A1",
         "Crear blacklists para 27 idiomas.\n"
         "IA genera borrador, Andrea valida con contexto cultural del contenido infantil.",
         "# de blacklists\nvalidadas y listas",
         "5 (Tier 1) en 30d\n27 en 60d",
         "Sem 3-8"],
        ["A2",
         "Revisar guiones ANTES de traduccion (Gate 1).\n"
         "Identificar modismos, expresiones coloquiales y humor que NO traduce bien.\n"
         "Ej: 'que padre' -> 'what a father'",
         "# expresiones\nflagged por episodio\n% errores evitados",
         "Checklist de\n20+ expresiones\npor cluster",
         "Sem 2-6"],
        ["A3",
         "Definir reglas culturales por cluster de idiomas.\n"
         "Que es aceptable en LATAM pero no en Medio Oriente,\n"
         "que humor funciona en Europa vs Asia.",
         "Documento de\nreglas por cluster\n(5 clusters min)",
         "1 cluster/semana",
         "Sem 4-10"],
        ["A4",
         "Monitorear contenido publicado.\n"
         "Reportar temas culturales o de lenguaje detectados post-publicacion.\n"
         "Meta: cero incidentes Cat. A (groserias, ofensas culturales) en Tier 1.",
         "0 incidentes Cat. A\nen Tier 1",
         "0 Cat. A en\nTier 1 siempre",
         "Continuo"],
    ]
    build_table(
        slide, Inches(0.3), Inches(1.45), Inches(12.73), Inches(5.05), rows,
        col_widths=[Inches(0.5), Inches(5.0), Inches(2.43), Inches(2.3), Inches(2.5)],
        font_size=Pt(10))

    add_highlight_box(slide, Inches(0.3), Inches(6.6), Inches(12.73), Inches(0.7),
                      "Andrea es el primer filtro: si el guion tiene problemas culturales, "
                      "se multiplican x27 idiomas.",
                      font_size=Pt(16))

    add_slide_number(slide, 18)


def slide_19_smart_alan(prs):
    """SLIDE 19: SMART — Alan (Optimizacion Post-Produccion)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_GRAY)
    add_slide_header(slide, "Alan \u2014 Optimizacion Post-Produccion Multi-Idioma",
                     "Eficiencia del equipo de traduccion, gaps, coordinacion")

    rows = [
        ["#", "Objetivo", "Como se mide", "Meta", "Plazo"],
        ["L1",
         "Reducir tiempo de correccion manual de Saul/Ivan.\n"
         "Identificar patrones repetitivos y eliminarlos con procesos o herramientas.",
         "Horas de correccion\npor episodio",
         "-20% en 60d\n-40% en 90d",
         "Sem 3-12"],
        ["L2",
         "Identificar gaps, silencios y problemas de timing en traducciones Tier 1.\n"
         "Escuchar muestras de cada idioma Tier 1 y reportar donde falla el audio.",
         "# gaps detectados\npor episodio\n% resueltos",
         "100% Tier 1\nrevisado cada\nepisodio",
         "Sem 4-10"],
        ["L3",
         "Coordinarse con Ramon para automatizar deteccion de gaps.\n"
         "Definir que problemas pueden detectarse con herramientas vs oido humano.",
         "# tipos de gap\nautomatizados vs\nmanuales",
         "50% deteccion\nautomatica en\nTier 1",
         "Sem 6-12"],
        ["L4",
         "Medir y reportar tasa de aprobacion a la primera (FTR).\n"
         "Cuantos episodios pasan QA sin necesidad de re-trabajo.",
         "FTR % semanal\n(First Time Right)",
         ">40% en 60d\n>60% en 90d",
         "Sem 3-12"],
    ]
    build_table(
        slide, Inches(0.3), Inches(1.45), Inches(12.73), Inches(5.05), rows,
        col_widths=[Inches(0.5), Inches(5.0), Inches(2.43), Inches(2.3), Inches(2.5)],
        font_size=Pt(10))

    add_highlight_box(slide, Inches(0.3), Inches(6.6), Inches(12.73), Inches(0.7),
                      "Alan optimiza personas + se apoya en Ramon para lo tecnico. "
                      "Objetivo: menos retrabajo, mas eficiencia en Saul/Ivan.",
                      font_size=Pt(16))

    add_slide_number(slide, 19)


def slide_20_smart_ramon(prs):
    """SLIDE 20: SMART — Ramon (Soporte Tecnico)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_GRAY)
    add_slide_header(slide, "Ramon \u2014 Soporte Tecnico del Pipeline",
                     "Herramientas, automatizacion, infraestructura")

    rows = [
        ["#", "Objetivo", "Como se mide", "Meta", "Plazo"],
        ["R1",
         "Configurar y validar pipeline end-to-end con 1 episodio piloto exitoso.\n"
         "Flujo Guion->ElevenLabs->Audio funcione sin errores manuales.",
         "Pipeline funcional\ncon 1 episodio\npiloto exitoso",
         "1 episodio\ncompleto sin\nerrores",
         "Sem 1-3"],
        ["R2",
         "Implementar WER baseline para 27 idiomas + forced alignment.\n"
         "Medir tasa de error de traduccion por idioma como punto de partida.",
         "WER medido\npor idioma\n(27 idiomas)",
         "WER baseline\npara 27 idiomas",
         "Sem 3-6"],
        ["R3",
         "Integrar blacklists de Andrea al Gate 1 automatico.\n"
         "Que las listas se apliquen antes de enviar a ElevenLabs.",
         "# idiomas con\nblacklist integrada\nen Gate 1",
         "5 (Tier 1)\nen sem 6\n27 en sem 10",
         "Sem 4-10"],
        ["R4",
         "Soporte tecnico al equipo (<24h para problemas criticos).\n"
         "Resolver problemas de ElevenLabs, APIs, formatos de audio.",
         "Tiempo de\nresolucion\npromedio",
         "<24h para\nproblemas\ncriticos",
         "Continuo"],
    ]
    build_table(
        slide, Inches(0.3), Inches(1.45), Inches(12.73), Inches(5.05), rows,
        col_widths=[Inches(0.5), Inches(5.0), Inches(2.43), Inches(2.3), Inches(2.5)],
        font_size=Pt(10))

    add_highlight_box(slide, Inches(0.3), Inches(6.6), Inches(12.73), Inches(0.7),
                      "Ramon es el puente entre lo que el equipo necesita "
                      "y lo que la tecnologia puede resolver.",
                      font_size=Pt(16))

    add_slide_number(slide, 20)


def slide_21_smart_daniel(prs):
    """SLIDE 21: SMART — Daniel (Deep Research + Habilitacion IA)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_GRAY)
    add_slide_header(slide, "Daniel \u2014 Deep Research y Habilitacion IA",
                     "Investigaciones profundas, herramientas AI, apoyo al equipo")

    rows = [
        ["#", "Objetivo", "Como se mide", "Meta", "Plazo"],
        ["D1",
         "Ejecutar 4 investigaciones profundas (Deep Research) con multiples IAs.\n"
         "Cada DR genera recomendaciones accionables para el equipo.",
         "4 DRs ejecutados\ncon outputs\nvalidados",
         "DR01+DR02 sem 3\nDR03+DR04 sem 6",
         "Sem 1-6"],
        ["D2",
         "Traducir hallazgos de los DRs en acciones para Andrea/Alan/Ramon.\n"
         "Ej: resultados DR02 -> reglas culturales que Andrea aplica.",
         "# recomendaciones\nimplementadas\npor el equipo",
         ">80% de\nrecomendaciones\nadoptadas",
         "Sem 4-10"],
        ["D3",
         "Configurar 3 jueces IA operativos para el equipo.\n"
         "Multi-model QA (Claude + GPT + Gemini), dashboards, reportes automaticos.",
         "Herramientas\noperativas y\nusables por equipo",
         "3 jueces IA\noperativos\nen sem 8",
         "Sem 4-10"],
        ["D4",
         "Calcular ROI por idioma con datos reales de YouTube Analytics.\n"
         "Cruzar costos de dubbing con revenue para decidir que idiomas mantener/pausar.",
         "ROI calculado\npor idioma",
         "5 (Tier 1) sem 6\n27 idiomas sem 10",
         "Sem 4-10"],
    ]
    build_table(
        slide, Inches(0.3), Inches(1.45), Inches(12.73), Inches(5.05), rows,
        col_widths=[Inches(0.5), Inches(5.0), Inches(2.43), Inches(2.3), Inches(2.5)],
        font_size=Pt(10))

    add_highlight_box(slide, Inches(0.3), Inches(6.6), Inches(12.73), Inches(0.7),
                      "Daniel investiga con IA, traduce hallazgos a acciones "
                      "y habilita al equipo con herramientas.",
                      font_size=Pt(16))

    add_slide_number(slide, 21)


def slide_22_deep_research(prs):
    """SLIDE 22: Las 4 Investigaciones Profundas (Deep Research)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_GRAY)
    add_slide_header(slide, "Las 4 Investigaciones Profundas \u2014 Deep Research",
                     "Que es DR: se le hace la misma pregunta a 5 IAs diferentes y se sintetizan las respuestas")

    add_textbox(slide, Inches(0.5), Inches(1.45), Inches(12.3), Inches(0.55),
                "Cada DR usa ChatGPT, Gemini, Claude, Perplexity y mas — se sintetizan para "
                "obtener vision completa sin sesgo de una sola IA.",
                font_size=Pt(13), bold=False, color=DARK_GRAY)

    drs = [
        ("DR01", "Cadena de Traduccion",
         "Conviene traducir directo ES->idioma o siempre via EN?\n"
         "PT-BR pierde solo 18% de AVD vs 35% de EN.\n"
         "Puede ser mejor ir directo para idiomas cercanos al espanol.",
         "Afecta a: Alan, Ramon",
         DARK_BLUE),
        ("DR02", "Adaptacion Cultural",
         "Framework de reglas culturales por cluster de idiomas.\n"
         "Que blacklists necesitamos, que humor no traduce,\n"
         "que es ofensivo en cada region (ej: cerdo en arabe).",
         "Afecta a: Andrea, Alan",
         GREEN),
        ("DR03", "Calidad TTS por Idioma",
         "Que tan bien suena ElevenLabs en cada idioma?\n"
         "CJK pierde 49% de audiencia.\n"
         "Es el TTS? la traduccion? la cultura? Necesitamos datos.",
         "Afecta a: Ramon, Alan",
         ORANGE),
        ("DR04", "QA Multi-IA Empresarial",
         "Como Netflix/Disney hacen QA multilenguaje con IA.\n"
         "Patrones de multi-model consensus, escalas MQM,\n"
         "pipelines de validacion cultural automatizada.",
         "Afecta a: Daniel, Ramon",
         LIGHT_BLUE),
    ]

    card_w = Inches(5.9)
    card_h = Inches(2.15)
    col1_x = Inches(0.4)
    col2_x = Inches(6.8)

    for i, (dr_id, title, desc, affects, color) in enumerate(drs):
        x = col1_x if i % 2 == 0 else col2_x
        y = Inches(2.1) + (i // 2) * Inches(2.4)

        card = add_shape_bg(slide, x, y, card_w, card_h, WHITE)
        add_shape_bg(slide, x, y, Inches(0.1), card_h, color)

        badge = add_shape_bg(slide, x + Inches(0.2), y + Inches(0.15),
                             Inches(0.8), Inches(0.4), color)
        set_text(badge.text_frame, dr_id, font_size=Pt(11), bold=True,
                 color=WHITE, alignment=PP_ALIGN.CENTER)

        add_textbox(slide, x + Inches(1.1), y + Inches(0.1), Inches(4.6),
                    Inches(0.4), title, font_size=Pt(15), bold=True,
                    color=DARK_BLUE)

        add_textbox(slide, x + Inches(0.25), y + Inches(0.6), Inches(5.45),
                    Inches(1.2), desc, font_size=Pt(10), bold=False,
                    color=DARK_GRAY)

        add_textbox(slide, x + Inches(0.25), y + Inches(1.8), Inches(5.45),
                    Inches(0.3), affects, font_size=Pt(10), bold=True,
                    color=color)

    add_slide_number(slide, 22)


def slide_23_acciones(prs):
    """SLIDE 23: Acciones esta Semana."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_GRAY)
    add_slide_header(slide, "Acciones esta Semana",
                     "Cada quien tiene una tarea concreta de bajo esfuerzo / alto impacto")

    rows = [
        ["Persona", "Accion", "Tiempo"],
        ["Andrea",
         "Listar 10 expresiones/modismos criticos del guion que NO traducen bien",
         "30 min"],
        ["Alan",
         "Agendar sesion de alineacion + revisar metricas de eficiencia actuales",
         "30 min"],
        ["Ramon",
         "Validar catalogo de voces (manifest.json) y reportar estado del pipeline",
         "1 hora"],
        ["Daniel",
         "Fase 0: verificar pipeline end-to-end + fix bugs criticos P0/P1",
         "Esta semana"],
        ["Saul / Ivan",
         "Confirmar disponibilidad para piloto de 1 episodio completo",
         "10 min"],
    ]
    build_table(
        slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(4.2), rows,
        col_widths=[Inches(2.0), Inches(8.3), Inches(2.0)],
        font_size=Pt(14))

    add_textbox(slide, Inches(0.5), Inches(5.85), Inches(12.3), Inches(0.45),
                "Tiempo total del equipo: ~3.5 horas | "
                "Impacto: desbloquea todo el proyecto",
                font_size=Pt(16), bold=True, color=ORANGE,
                alignment=PP_ALIGN.CENTER)

    add_highlight_box(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.85),
                      "Estas acciones son el unico prerequisito para arrancar el piloto en Semana 2.",
                      font_size=Pt(18))

    add_slide_number(slide, 23)


def slide_24_decisiones(prs):
    """SLIDE 24: Decisiones Pendientes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_GRAY)
    add_slide_header(slide, "Decisiones Pendientes",
                     "El equipo necesita resolver estas 5 preguntas")

    decisions = [
        ("1", "Cuantos idiomas en el piloto?",
         "Opciones: 5 (Tier 1), 10 (Tier 1+2), o los 27 completos desde el inicio",
         DARK_BLUE),
        ("2", "Quien es dueno permanente de las blacklists?",
         "Andrea las crea — pero quien las actualiza y aprueba a largo plazo?",
         ORANGE),
        ("3", "Pausar idiomas de bajo ROI?",
         "Tamil (30K views, AVD 1:58), Filipino (46K views), Chino (65K, revenue cero)",
         RED),
        ("4", "Migrar a Enterprise de ElevenLabs?",
         "Mas capacidad, menor costo por caracter — depende del volumen mensual",
         LIGHT_BLUE),
        ("5", "Que tan estricto debe ser el sistema?",
         "Bloquear automaticamente (arriesga retrasos) vs alertar y dejar pasar (arriesga errores)",
         GREEN),
    ]

    for i, (num, question, detail, color) in enumerate(decisions):
        y = Inches(1.5) + i * Inches(1.0)

        num_badge = add_shape_bg(slide, Inches(0.4), y + Inches(0.1),
                                 Inches(0.5), Inches(0.65), color)
        set_text(num_badge.text_frame, num, font_size=Pt(18), bold=True,
                 color=WHITE, alignment=PP_ALIGN.CENTER)

        add_textbox(slide, Inches(1.05), y + Inches(0.05), Inches(11.5),
                    Inches(0.4), question, font_size=Pt(17), bold=True,
                    color=DARK_BLUE)
        add_textbox(slide, Inches(1.05), y + Inches(0.45), Inches(11.5),
                    Inches(0.45), detail, font_size=Pt(13), bold=False,
                    color=DARK_GRAY)

    add_slide_number(slide, 24)


def slide_25_siguiente_paso(prs):
    """SLIDE 25: Siguiente Paso."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_GRAY)
    add_slide_header(slide, "Siguiente Paso",
                     "Plan de las proximas 4 semanas")

    steps = [
        ("HOY",
         "Alinear vision del equipo + asignar objetivos SMART",
         DARK_BLUE),
        ("ESTA SEMANA",
         "Andrea: blacklist Tier 1  |  Alan: baseline tiempos  |  Ramon: validar pipeline",
         ORANGE),
        ("SEMANA 2",
         "Piloto con 1 episodio real (27 idiomas) — medir todo por primera vez",
         LIGHT_BLUE),
        ("SEMANA 3",
         "Daniel: DR01+DR02 listos, primeras recomendaciones para el equipo",
         GREEN),
        ("MES 1",
         "Metricas funcionando para Tier 1, primeros FTR medidos y publicados",
         DARK_BLUE),
    ]

    for i, (when, what, color) in enumerate(steps):
        y = Inches(1.55) + i * Inches(1.05)

        badge = add_shape_bg(slide, Inches(0.5), y, Inches(2.5),
                             Inches(0.75), color)
        set_text(badge.text_frame, when, font_size=Pt(16), bold=True,
                 color=WHITE, alignment=PP_ALIGN.CENTER)

        add_textbox(slide, Inches(3.1), y, Inches(0.5), Inches(0.75),
                    "\u2192", font_size=Pt(24), bold=True, color=color,
                    alignment=PP_ALIGN.CENTER)

        add_textbox(slide, Inches(3.7), y + Inches(0.1), Inches(9.1),
                    Inches(0.6), what, font_size=Pt(16), bold=False,
                    color=DARK_GRAY)

        if i < len(steps) - 1:
            line_y = y + Inches(0.75)
            add_shape_bg(slide, Inches(1.7), line_y, Inches(0.08),
                         Inches(0.3), MEDIUM_GRAY)

    add_slide_number(slide, 25)


def slide_26_gracias(prs):
    """SLIDE 26: Gracias / Preguntas."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BLUE)

    add_shape_bg(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.15), ORANGE)

    add_textbox(slide, Inches(3.5), Inches(1.2), Inches(6.3), Inches(1.6),
                "Preguntas?", font_size=Pt(52), bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(2.5), Inches(3.1), Inches(8.3), Inches(0.6),
                "Contacto: Daniel Garza", font_size=Pt(22), bold=False,
                color=ORANGE, alignment=PP_ALIGN.CENTER)

    add_shape_bg(slide, Inches(4.0), Inches(3.9), Inches(5.3), Inches(0.03),
                 ORANGE)

    add_textbox(slide, Inches(1.0), Inches(4.3), Inches(11.3), Inches(2.0),
                "Este documento resume 79 archivos de analisis, 4 perspectivas AI,\n"
                "datos de YouTube Analytics 2025, 4 Deep Research en curso\n"
                "y 129 docs de API ElevenLabs.",
                font_size=Pt(16), bold=False,
                color=RGBColor(0x99, 0xAA, 0xCC), alignment=PP_ALIGN.CENTER)

    add_shape_bg(slide, Inches(0), Inches(7.35), SLIDE_WIDTH, Inches(0.15), ORANGE)

    add_slide_number(slide, 26)


# ---------------------------------------------------------------------------
# MAIN
# ---------------------------------------------------------------------------

def main():
    """Generate the full presentation — v6, 26 slides."""
    prs = Presentation()

    # Set 16:9 slide dimensions
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # PART 1: DONDE ESTAMOS (slides 1-5)
    slide_01_title(prs)              # 1 — Title
    slide_02_agenda(prs)             # 2 — Agenda
    slide_03_panorama(prs)           # 3 — Panorama datos reales
    slide_04_tabla_idiomas(prs)      # 4 — Tabla ingreso por idioma
    slide_05_perdida_audiencia(prs)  # 5 — AVD bars / perdida audiencia

    # PART 2: LOS 6 PROBLEMAS (slides 6-12)
    slide_06_problema1(prs)          # 6 — Telefono descompuesto
    slide_07_problema2(prs)          # 7 — Zero metricas
    slide_08_problema3(prs)          # 8 — Blacklists solo 13 palabras
    slide_09_problema4(prs)          # 9 — Desconexion cultural
    slide_10_problema5(prs)          # 10 — Dos sistemas desconectados
    slide_11_problema6(prs)          # 11 — Inversion a ciegas
    slide_12_cerdo_robot(prs)        # 12 — Caso del cerdo robot

    # PART 3: LA SOLUCION (slides 13-16)
    slide_13_buena_noticia(prs)      # 13 — Buena noticia 60%
    slide_14_solucion(prs)           # 14 — 4 Gates de calidad
    slide_15_costo(prs)              # 15 — Costo del sistema
    slide_16_tiers(prs)              # 16 — Prioridad por tiers

    # PART 4: SMART Y DEEP RESEARCH (slides 17-22)
    slide_17_smart_overview(prs)     # 17 — SMART overview
    slide_18_smart_andrea(prs)       # 18 — Andrea SMART
    slide_19_smart_alan(prs)         # 19 — Alan SMART
    slide_20_smart_ramon(prs)        # 20 — Ramon SMART
    slide_21_smart_daniel(prs)       # 21 — Daniel SMART
    slide_22_deep_research(prs)      # 22 — 4 Deep Research

    # PART 5: CIERRE (slides 23-26)
    slide_23_acciones(prs)           # 23 — Acciones esta semana
    slide_24_decisiones(prs)         # 24 — Decisiones pendientes
    slide_25_siguiente_paso(prs)     # 25 — Siguiente paso
    slide_26_gracias(prs)            # 26 — Gracias / Preguntas

    # Save
    prs.save(str(OUTPUT_PATH))
    print(f"Presentation saved to: {OUTPUT_PATH}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
